from __future__ import annotations

from datetime import datetime, timezone
import hashlib
import os
import re
import logging
from io import BytesIO
from typing import Any
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader

from azure.core.credentials import AzureKeyCredential
from azure.search.documents import SearchClient
from azure.search.documents.indexes import SearchIndexClient

from backend.storage.blob_storage import BlobStorageClient
from backend.knowledge.embeddings import generate_embedding


class KnowledgeSearchIndex:
    def __init__(self, endpoint: str, index_name: str, admin_key: str) -> None:
        self._endpoint = endpoint
        self._index_name = index_name
        self._credential = AzureKeyCredential(admin_key)
        self._search_client = SearchClient(
            endpoint=endpoint,
            index_name=index_name,
            credential=self._credential,
        )
        self._index_client = SearchIndexClient(
            endpoint=endpoint,
            credential=self._credential,
        )

    @property
    def index_name(self) -> str:
        return self._index_name

    def upload_documents(self, documents: list[dict[str, Any]]) -> list:
        if not isinstance(documents, list):
            raise TypeError(
                f"upload_documents expects list[dict], got {type(documents).__name__}"
            )
        return self._search_client.upload_documents(documents=documents)


class KnowledgeIngestionService:
    def __init__(
        self,
        blob_client: BlobStorageClient,
        search_index: KnowledgeSearchIndex,
        prefix: str = "knowledge/",
    ) -> None:
        self._blob_client = blob_client
        self._prefix = prefix
        self._search_index = search_index
        self._logger = logging.getLogger("knowledge_ingestion")

    def delete_by_source(self, filename: str) -> int:
        """
        Delete all index documents whose source field matches filename.
        Returns the count of documents deleted.
        """
        deleted = 0
        try:
            results = self._search_index._search_client.search(
                search_text="*",
                filter=f"source eq '{filename}'",
                select=["doc_id"],
                top=1000,
            )
            doc_ids = [r["doc_id"] for r in results if r.get("doc_id")]
            if doc_ids:
                delete_batch = [{"doc_id": doc_id} for doc_id in doc_ids]
                self._search_index._search_client.delete_documents(
                    documents=delete_batch
                )
                deleted = len(doc_ids)
                self._logger.info(
                    "[KNOWLEDGE] Deleted %d existing chunks for '%s' before re-ingestion",
                    deleted,
                    filename,
                )
        except Exception as e:
            self._logger.warning(
                "[KNOWLEDGE] delete_by_source failed for '%s': %s", filename, e
            )
        return deleted

    def upload_document(self, filename: str, data: bytes, content_type: str) -> None:
        # Delete any existing index entries for this file before re-ingesting
        self.delete_by_source(filename)
        path = f"{self._prefix}{filename}"
        self._blob_client.upload_file(path, data, content_type, overwrite=True)
        text = self._extract_text(data, content_type, filename)

        base_doc_id = self._build_doc_id(filename)
        created_at = datetime.now(timezone.utc).isoformat()

        # STEP 1 — Build document_summary entry
        summary_text = text.strip()[:500]
        summary_id = base_doc_id + "_summary"
        summary_embedding = generate_embedding(summary_text)
        if not isinstance(summary_embedding, list) or len(summary_embedding) != 3072:
            raise ValueError(
                f"Invalid embedding length for summary of '{filename}': "
                f"expected 3072, got "
                f"{len(summary_embedding) if isinstance(summary_embedding, list) else 'non-list'}"
            )
        summary_doc: dict = {
            "doc_id": summary_id,
            "doc_type": "knowledge",
            "title": filename,
            "content_text": summary_text,
            "source": filename,
            "version": "1",
            "created_at": created_at,
            "chunk_type": "document_summary",
            "section_title": filename,
            "parent_section_id": "",
            "page_start": 0,
            "page_end": 0,
            "cosolve_phase": self._detect_cosolve_phase(text),
            "char_count": len(summary_text),
            "embedding": summary_embedding,
        }

        # STEP 2 — Split into sections
        sections = self._split_into_sections(text, filename)
        section_docs: list[dict] = []
        all_small_chunk_docs: list[dict] = []

        for idx, section in enumerate(sections):
            section_id = f"{base_doc_id}_sec_{idx}"
            cosolve_phase = self._detect_cosolve_phase(section["content"])

            section_embedding = generate_embedding(
                section["content"]
            )
            if (
                not isinstance(section_embedding, list)
                or len(section_embedding) != 3072
            ):
                raise ValueError(
                    f"Invalid embedding length for section {idx} of '{filename}': "
                    f"expected 3072, got "
                    f"{len(section_embedding) if isinstance(section_embedding, list) else 'non-list'}"
                )
            section_docs.append(
                {
                    "doc_id": section_id,
                    "doc_type": "knowledge",
                    "title": filename,
                    "content_text": section["content"],
                    "source": filename,
                    "version": "1",
                    "created_at": created_at,
                    "chunk_type": "section",
                    "section_title": section["section_title"],
                    "parent_section_id": "",
                    "page_start": section["page_start"],
                    "page_end": section["page_end"],
                    "cosolve_phase": cosolve_phase,
                    "char_count": len(section["content"]),
                    "embedding": section_embedding,
                }
            )

            small_chunk_dicts = self._build_small_chunks(
                section["content"],
                section_id,
                filename,
                section["section_title"],
                cosolve_phase,
                created_at,
            )
            for sc in small_chunk_dicts:
                sc_embedding = generate_embedding(
                    sc["content_text"]
                )
                if not isinstance(sc_embedding, list) or len(sc_embedding) != 3072:
                    raise ValueError(
                        f"Invalid embedding length for small chunk of '{filename}': "
                        f"expected 3072, got "
                        f"{len(sc_embedding) if isinstance(sc_embedding, list) else 'non-list'}"
                    )
                sc["embedding"] = sc_embedding
                all_small_chunk_docs.append(sc)

        # STEP 3 — Collect all documents and upload in one batch
        documents_to_upload = [summary_doc] + section_docs + all_small_chunk_docs
        self._search_index.upload_documents(documents_to_upload)

        # STEP 4 — Log summary
        total_small_chunks = len(all_small_chunk_docs)
        self._logger.info(
            "[KNOWLEDGE] '%s' → 1 summary + %d sections + %d small chunks",
            filename,
            len(sections),
            total_small_chunks,
        )
        print(
            f"[KNOWLEDGE] '{filename}' → 1 summary + {len(sections)} sections "
            f"+ {total_small_chunks} small chunks"
        )

    def delete_knowledge_blob(self, filename: str) -> None:
        """Delete an orphaned blob from storage (no index record required)."""
        path = f"{self._prefix}{filename}"
        self._blob_client.delete_file(path)
        self._logger.info("[KNOWLEDGE] deleted orphan blob: %s", path)
        print(f"[KNOWLEDGE] deleted orphan blob: {path}")

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    def _split_into_sections(self, text: str, filename: str) -> list[dict]:
        """Split document text into logical sections.

        STEP 1: heading-based splitting (3+ headings detected).
        STEP 2: fixed-size 2000-char / 200-char-overlap fallback.
        """
        lines = text.splitlines()

        # STEP 1 — Try heading detection
        heading_indices: list[int] = []
        for i, line in enumerate(lines):
            stripped = line.strip()
            if not stripped:
                continue
            is_heading = False
            if re.match(r"^[A-Z][A-Z0-9 \-/&:,\.]{2,79}$", stripped):
                is_heading = True
            elif re.match(r"^\d{1,2}(\.\d+)*[\.\s]+\S", stripped):
                is_heading = True
            elif len(stripped) < 60:
                next_line = lines[i + 1] if i + 1 < len(lines) else ""
                if next_line.strip() == "":
                    # Additional guards: reject if line contains sentence-ending
                    # punctuation mid-string or looks like body text
                    has_sentence_end = bool(re.search(r"[.!?]\s+\S", stripped))
                    starts_with_lower = stripped[0].islower() if stripped else False
                    is_conjunction_start = re.match(
                        r"^(the|a|an|this|that|these|those|if|when|where|"
                        r"as|by|for|in|on|or|and|but|to|of|with|from)\b",
                        stripped,
                        re.IGNORECASE,
                    )
                    if (
                        not has_sentence_end
                        and not starts_with_lower
                        and not is_conjunction_start
                    ):
                        is_heading = True
            if is_heading:
                heading_indices.append(i)

        if len(heading_indices) >= 3:
            raw_sections: list[dict] = []
            for n, hi in enumerate(heading_indices):
                end_line = (
                    heading_indices[n + 1]
                    if n + 1 < len(heading_indices)
                    else len(lines)
                )
                content = "\n".join(lines[hi:end_line]).strip()
                title = lines[hi].strip()
                if len(title) > 120:
                    title = title[:120].rsplit(" ", 1)[0]
                if len(content) < 200:
                    continue  # likely a TOC entry
                raw_sections.append(
                    {
                        "section_title": title,
                        "content": content,
                        "page_start": 0,
                        "page_end": 0,
                    }
                )

            # Sub-split sections exceeding 3000 chars at paragraph breaks
            sections: list[dict] = []
            for sec in raw_sections:
                if len(sec["content"]) <= 3000:
                    sections.append(sec)
                else:
                    content = sec["content"]
                    start = 0
                    part_n = 1
                    while start < len(content):
                        end = start + 3000
                        if end >= len(content):
                            sub = content[start:].strip()
                            if sub:
                                label = (
                                    sec["section_title"]
                                    if part_n == 1
                                    else f"{sec['section_title']} (continued)"
                                )
                                sections.append(
                                    {
                                        "section_title": label,
                                        "content": sub,
                                        "page_start": 0,
                                        "page_end": 0,
                                    }
                                )
                            break
                        split_at = content.rfind("\n\n", start, end)
                        if split_at <= start:
                            split_at = content.rfind("\n", start, end)
                        if split_at <= start:
                            split_at = end
                        sub = content[start:split_at].strip()
                        if sub:
                            label = (
                                sec["section_title"]
                                if part_n == 1
                                else f"{sec['section_title']} (continued)"
                            )
                            sections.append(
                                {
                                    "section_title": label,
                                    "content": sub,
                                    "page_start": 0,
                                    "page_end": 0,
                                }
                            )
                        # advance past the split point, skipping leading whitespace
                        start = split_at
                        while start < len(content) and content[start] in (
                            "\n",
                            " ",
                            "\r",
                        ):
                            start += 1
                        part_n += 1

            if sections:
                return sections

        # STEP 2 — Fallback: fixed-size 2000-char split with 200-char overlap
        fallback_sections: list[dict] = []
        start = 0
        n = 1
        while start < len(text):
            end = start + 2000
            if end >= len(text):
                chunk = text[start:].strip()
                if chunk:
                    fallback_sections.append(
                        {
                            "section_title": f"{filename} \u2014 Part {n}",
                            "content": chunk,
                            "page_start": 0,
                            "page_end": 0,
                        }
                    )
                break
            split_at = text.rfind("\n\n", start, end)
            if split_at <= start:
                split_at = text.rfind("\n", start, end)
            if split_at <= start:
                split_at = end
            chunk = text[start:split_at].strip()
            if chunk:
                fallback_sections.append(
                    {
                        "section_title": f"{filename} \u2014 Part {n}",
                        "content": chunk,
                        "page_start": 0,
                        "page_end": 0,
                    }
                )
            next_start = split_at - 200
            if next_start <= start:
                next_start = split_at
            start = next_start
            n += 1

        if not fallback_sections:
            fallback_sections = [
                {
                    "section_title": f"{filename} \u2014 Part 1",
                    "content": text,
                    "page_start": 0,
                    "page_end": 0,
                }
            ]
        return fallback_sections

    def _detect_cosolve_phase(self, text: str) -> str:
        """Tag text with the most relevant CoSolve reasoning phase.

        Returns one of: 'diagnose' | 'root_cause' | 'correct' | 'prevent' | 'general'
        """
        lower = text.lower()
        phase_keywords: dict[str, list[str]] = {
            "diagnose": [
                "symptom",
                "alarm",
                "temperature",
                "observed",
                "detected",
                "failure",
                "overheating",
                "reading",
                "telemetry",
                "monitoring",
                "measurement",
            ],
            "root_cause": [
                "cause",
                "root cause",
                "failure mechanism",
                "why",
                "analysis",
                "confirmed",
                "investigation",
                "factor",
                "contributed",
                "determined",
                "evidence",
            ],
            "correct": [
                "corrective action",
                "replacement",
                "revised",
                "updated",
                "repair",
                "replaced",
                "implemented",
                "fixed",
                "rework",
                "modification",
                "action taken",
            ],
            "prevent": [
                "prevention",
                "systemic",
                "recurrence",
                "policy",
                "procedure",
                "training",
                "audit",
                "qualification",
                "specification",
                "standard",
                "requirement",
                "incoming inspection",
                "supplier",
            ],
        }
        scores: dict[str, int] = {phase: 0 for phase in phase_keywords}
        for phase, keywords in phase_keywords.items():
            for kw in keywords:
                if kw in lower:
                    scores[phase] += 1
        best_phase = max(scores, key=lambda p: scores[p])
        if scores[best_phase] == 0:
            return "general"
        # Check for tie
        top_score = scores[best_phase]
        tied = [p for p, s in scores.items() if s == top_score]
        if len(tied) > 1:
            return "general"
        return best_phase

    def _build_small_chunks(
        self,
        section_content: str,
        section_id: str,
        source: str,
        section_title: str,
        cosolve_phase: str,
        created_at: str,
    ) -> list[dict]:
        """Create small_chunk index documents (500 chars, 100-char overlap) from a section."""
        small_chunks: list[dict] = []
        text = section_content
        start = 0
        idx = 0
        while start < len(text):
            end = start + 500
            if end >= len(text):
                chunk_text = text[start:].strip()
                if chunk_text:
                    small_chunks.append(
                        {
                            "doc_id": f"{section_id}_sc_{idx}",
                            "doc_type": "knowledge",
                            "title": source,
                            "content_text": chunk_text,
                            "source": source,
                            "version": "1",
                            "created_at": created_at,
                            "chunk_type": "small_chunk",
                            "section_title": section_title,
                            "parent_section_id": section_id,
                            "page_start": 0,
                            "page_end": 0,
                            "cosolve_phase": cosolve_phase,
                            "char_count": len(chunk_text),
                            "embedding": None,
                        }
                    )
                break
            split_at = text.rfind(" ", start, end)
            if split_at <= start:
                split_at = end
            chunk_text = text[start:split_at].strip()
            if chunk_text:
                small_chunks.append(
                    {
                        "doc_id": f"{section_id}_sc_{idx}",
                        "doc_type": "knowledge",
                        "title": source,
                        "content_text": chunk_text,
                        "source": source,
                        "version": "1",
                        "created_at": created_at,
                        "chunk_type": "small_chunk",
                        "section_title": section_title,
                        "parent_section_id": section_id,
                        "page_start": 0,
                        "page_end": 0,
                        "cosolve_phase": cosolve_phase,
                        "char_count": len(chunk_text),
                        "embedding": None,
                    }
                )
            next_start = split_at - 100
            if next_start <= start:
                next_start = split_at
            start = next_start
            idx += 1
        return small_chunks

    def _build_doc_id(self, filename: str) -> str:
        digest = hashlib.sha1(filename.encode("utf-8")).hexdigest()
        return digest

    def _extract_text(self, data: bytes, content_type: str, filename: str) -> str:
        if not data:
            raise ValueError("Extracted text is empty")

        ext = os.path.splitext(filename or "")[1].lower()

        if ext == ".txt":
            text = data.decode("utf-8", errors="ignore").strip()
        elif ext == ".docx":
            doc = Document(BytesIO(data))
            parts: list[str] = []
            # 1. Top-level paragraphs
            for p in doc.paragraphs:
                if isinstance(p.text, str) and p.text:
                    parts.append(p.text)
            # 2. Table cells (all rows, all cells, all paragraphs within each cell)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for p in cell.paragraphs:
                            if isinstance(p.text, str) and p.text:
                                parts.append(p.text)
            # 3. Headers and footers from each section
            for section in doc.sections:
                for p in section.header.paragraphs:
                    if isinstance(p.text, str) and p.text:
                        parts.append(p.text)
                for p in section.footer.paragraphs:
                    if isinstance(p.text, str) and p.text:
                        parts.append(p.text)
            text = "\n".join(parts).strip()
        elif ext == ".pdf":
            # PDF extraction is handled by a dedicated fallback chain.
            return self._extract_pdf_text(data, filename)
        elif ext == ".pptx":
            presentation = Presentation(BytesIO(data))
            parts: list[str] = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    value = getattr(shape, "text", "")
                    if isinstance(value, str) and value:
                        parts.append(value)
            text = "\n".join(parts).strip()
        else:
            raise ValueError("Unsupported file format")

        if not text:
            raise ValueError("Extracted text is empty")
        return text

    def _extract_pdf_text(self, data: bytes, filename: str) -> str:
        """Try PyPDF2 → pypdf → pdfplumber, returning a placeholder if all fail."""

        # --- Extractor 1: PyPDF2 (primary, already a dependency) ---
        try:
            reader = PdfReader(BytesIO(data))
            parts = [page.extract_text() or "" for page in reader.pages]
            text = "\n".join(p for p in parts if p).strip()
            if text:
                self._logger.debug(
                    "[KNOWLEDGE] PDF via PyPDF2: len=%d  file=%r", len(text), filename
                )
                return text
            self._logger.debug(
                "[KNOWLEDGE] PyPDF2 empty for %r, trying pypdf", filename
            )
        except Exception as exc:
            self._logger.warning(
                "[KNOWLEDGE] PyPDF2 failed for %r: %s: %s",
                filename,
                type(exc).__name__,
                exc,
            )

        # --- Extractor 2: pypdf 6.x (modern fork, already in requirements) ---
        try:
            from pypdf import PdfReader as PypdfReader

            reader2 = PypdfReader(BytesIO(data))
            parts2 = [page.extract_text() or "" for page in reader2.pages]
            text2 = "\n".join(p for p in parts2 if p).strip()
            if text2:
                self._logger.debug(
                    "[KNOWLEDGE] PDF via pypdf: len=%d  file=%r", len(text2), filename
                )
                return text2
            self._logger.debug(
                "[KNOWLEDGE] pypdf empty for %r, trying pdfplumber", filename
            )
        except Exception as exc:
            self._logger.warning(
                "[KNOWLEDGE] pypdf failed for %r: %s: %s",
                filename,
                type(exc).__name__,
                exc,
            )

        # --- Extractor 3: pdfplumber (handles more complex layouts) ---
        try:
            import pdfplumber

            with pdfplumber.open(BytesIO(data)) as pdf:
                parts3 = [pg.extract_text() or "" for pg in pdf.pages]
            text3 = "\n".join(p for p in parts3 if p).strip()
            if text3:
                self._logger.debug(
                    "[KNOWLEDGE] PDF via pdfplumber: len=%d  file=%r",
                    len(text3),
                    filename,
                )
                return text3
            self._logger.debug("[KNOWLEDGE] pdfplumber empty for %r", filename)
        except Exception as exc:
            self._logger.warning(
                "[KNOWLEDGE] pdfplumber failed for %r: %s: %s",
                filename,
                type(exc).__name__,
                exc,
            )

        # --- All extractors exhausted (likely a scanned / image-only PDF) ---
        self._logger.warning(
            "[KNOWLEDGE] All PDF extractors returned empty for %r "
            "(scanned/image-only PDF?). Indexing with placeholder so blob is not orphaned.",
            filename,
        )
        print(
            f"[KNOWLEDGE] WARNING: no extractable text in PDF {filename!r}; "
            "indexing with '[No extractable text]' placeholder"
        )
        return "[No extractable text]"


__all__ = [
    "KnowledgeIngestionService",
    "KnowledgeSearchIndex",
]
