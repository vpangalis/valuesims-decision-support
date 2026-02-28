from __future__ import annotations

from datetime import datetime, timezone
import hashlib
import os
import logging
from io import BytesIO
from typing import Any
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader

from azure.core.credentials import AzureKeyCredential
from azure.search.documents import SearchClient
from azure.search.documents.indexes import SearchIndexClient

from backend.infra.blob_storage import BlobStorageClient
from backend.infra.embeddings import EmbeddingClient


class KnowledgeSearchIndex:
    def __init__(self, endpoint: str, index_name: str, admin_key: str) -> None:
        self._endpoint = endpoint
        self._index_name = index_name
        self._credential = AzureKeyCredential(admin_key)
        self._search_client = SearchClient(
            endpoint=endpoint,
            index_name=index_name,
            credential=self._credential,
        )
        self._index_client = SearchIndexClient(
            endpoint=endpoint,
            credential=self._credential,
        )

    @property
    def index_name(self) -> str:
        return self._index_name

    def upload_documents(self, documents: list[dict[str, Any]]) -> list:
        if not isinstance(documents, list):
            raise TypeError(
                f"upload_documents expects list[dict], got {type(documents).__name__}"
            )
        return self._search_client.upload_documents(documents=documents)


class KnowledgeIngestionService:
    def __init__(
        self,
        blob_client: BlobStorageClient,
        embedding_client: EmbeddingClient,
        search_index: KnowledgeSearchIndex,
        prefix: str = "knowledge/",
    ) -> None:
        self._blob_client = blob_client
        self._prefix = prefix
        self._embedding_client = embedding_client
        self._search_index = search_index
        self._logger = logging.getLogger("knowledge_ingestion")

    # Maximum characters per chunk (≈ 3 000 tokens × 4 chars/token, safely under 8 192 token limit)
    _CHUNK_MAX_CHARS: int = 12_000

    def upload_document(self, filename: str, data: bytes, content_type: str) -> None:
        path = f"{self._prefix}{filename}"
        self._blob_client.upload_file(path, data, content_type, overwrite=True)
        text = self._extract_text(data, content_type, filename)

        chunks = self._chunk_text(text)
        total_chunks = len(chunks)
        base_doc_id = self._build_doc_id(filename)
        created_at = datetime.now(timezone.utc).isoformat()

        self._logger.info(
            "[KNOWLEDGE] '%s' → %d chunk(s), total chars=%d",
            filename,
            total_chunks,
            len(text),
        )
        print(
            f"[KNOWLEDGE] '{filename}' → {total_chunks} chunk(s), "
            f"total_chars={len(text)}"
        )

        documents_to_upload: list[dict] = []
        for idx, chunk in enumerate(chunks):
            doc_id = base_doc_id if total_chunks == 1 else f"{base_doc_id}_chunk_{idx}"

            embedding = self._embedding_client.generate_embedding(chunk)
            if not isinstance(embedding, list) or len(embedding) != 3072:
                raise ValueError(
                    f"Invalid embedding length for knowledge ingestion chunk {idx}: "
                    f"expected 3072, got "
                    f"{len(embedding) if isinstance(embedding, list) else 'non-list'}"
                )

            self._logger.info(
                "[KNOWLEDGE] chunk %d/%d — doc_id=%s  chars=%d  embedding_len=%d",
                idx,
                total_chunks - 1,
                doc_id,
                len(chunk),
                len(embedding),
            )
            print(
                f"[KNOWLEDGE] chunk {idx}/{total_chunks - 1} — doc_id={doc_id}  "
                f"chars={len(chunk)}  embedding_len={len(embedding)}"
            )

            documents_to_upload.append(
                {
                    "doc_id": doc_id,
                    "doc_type": "knowledge",
                    "title": filename,
                    "content_text": chunk,
                    "source": filename,
                    "version": "1",
                    "created_at": created_at,
                    "embedding": embedding,
                }
            )

        self._search_index.upload_documents(documents_to_upload)
        self._logger.info(
            "[KNOWLEDGE] uploaded %d chunk document(s) for '%s'",
            len(documents_to_upload),
            filename,
        )

    def delete_knowledge_blob(self, filename: str) -> None:
        """Delete an orphaned blob from storage (no index record required)."""
        path = f"{self._prefix}{filename}"
        self._blob_client.delete_file(path)
        self._logger.info("[KNOWLEDGE] deleted orphan blob: %s", path)
        print(f"[KNOWLEDGE] deleted orphan blob: {path}")

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    def _chunk_text(self, text: str) -> list[str]:
        """Split *text* into chunks of at most _CHUNK_MAX_CHARS characters.

        Splits on whitespace boundaries to avoid cutting mid-word.
        Returns a list with at least one element (even if text is empty).
        """
        max_chars = self._CHUNK_MAX_CHARS
        if len(text) <= max_chars:
            return [text]

        chunks: list[str] = []
        start = 0
        while start < len(text):
            end = start + max_chars
            if end >= len(text):
                chunks.append(text[start:])
                break
            # Walk back to the last whitespace so we don't split mid-word
            split_at = text.rfind(" ", start, end)
            if split_at <= start:  # no whitespace found — hard split
                split_at = end
            chunks.append(text[start:split_at])
            start = split_at + 1  # skip the space
        return chunks

    def _build_doc_id(self, filename: str) -> str:
        digest = hashlib.sha1(filename.encode("utf-8")).hexdigest()
        return digest

    def _extract_text(self, data: bytes, content_type: str, filename: str) -> str:
        if not data:
            raise ValueError("Extracted text is empty")

        ext = os.path.splitext(filename or "")[1].lower()

        if ext == ".txt":
            text = data.decode("utf-8", errors="ignore").strip()
        elif ext == ".docx":
            doc = Document(BytesIO(data))
            parts: list[str] = []
            # 1. Top-level paragraphs
            for p in doc.paragraphs:
                if isinstance(p.text, str) and p.text:
                    parts.append(p.text)
            # 2. Table cells (all rows, all cells, all paragraphs within each cell)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for p in cell.paragraphs:
                            if isinstance(p.text, str) and p.text:
                                parts.append(p.text)
            # 3. Headers and footers from each section
            for section in doc.sections:
                for p in section.header.paragraphs:
                    if isinstance(p.text, str) and p.text:
                        parts.append(p.text)
                for p in section.footer.paragraphs:
                    if isinstance(p.text, str) and p.text:
                        parts.append(p.text)
            text = "\n".join(parts).strip()
        elif ext == ".pdf":
            # PDF extraction is handled by a dedicated fallback chain.
            return self._extract_pdf_text(data, filename)
        elif ext == ".pptx":
            presentation = Presentation(BytesIO(data))
            parts: list[str] = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    value = getattr(shape, "text", "")
                    if isinstance(value, str) and value:
                        parts.append(value)
            text = "\n".join(parts).strip()
        else:
            raise ValueError("Unsupported file format")

        if not text:
            raise ValueError("Extracted text is empty")
        return text

    def _extract_pdf_text(self, data: bytes, filename: str) -> str:
        """Try PyPDF2 → pypdf → pdfplumber, returning a placeholder if all fail."""

        # --- Extractor 1: PyPDF2 (primary, already a dependency) ---
        try:
            reader = PdfReader(BytesIO(data))
            parts = [page.extract_text() or "" for page in reader.pages]
            text = "\n".join(p for p in parts if p).strip()
            if text:
                self._logger.debug(
                    "[KNOWLEDGE] PDF via PyPDF2: len=%d  file=%r", len(text), filename
                )
                return text
            self._logger.debug(
                "[KNOWLEDGE] PyPDF2 empty for %r, trying pypdf", filename
            )
        except Exception as exc:
            self._logger.warning(
                "[KNOWLEDGE] PyPDF2 failed for %r: %s: %s",
                filename,
                type(exc).__name__,
                exc,
            )

        # --- Extractor 2: pypdf 6.x (modern fork, already in requirements) ---
        try:
            from pypdf import PdfReader as PypdfReader

            reader2 = PypdfReader(BytesIO(data))
            parts2 = [page.extract_text() or "" for page in reader2.pages]
            text2 = "\n".join(p for p in parts2 if p).strip()
            if text2:
                self._logger.debug(
                    "[KNOWLEDGE] PDF via pypdf: len=%d  file=%r", len(text2), filename
                )
                return text2
            self._logger.debug(
                "[KNOWLEDGE] pypdf empty for %r, trying pdfplumber", filename
            )
        except Exception as exc:
            self._logger.warning(
                "[KNOWLEDGE] pypdf failed for %r: %s: %s",
                filename,
                type(exc).__name__,
                exc,
            )

        # --- Extractor 3: pdfplumber (handles more complex layouts) ---
        try:
            import pdfplumber

            with pdfplumber.open(BytesIO(data)) as pdf:
                parts3 = [pg.extract_text() or "" for pg in pdf.pages]
            text3 = "\n".join(p for p in parts3 if p).strip()
            if text3:
                self._logger.debug(
                    "[KNOWLEDGE] PDF via pdfplumber: len=%d  file=%r",
                    len(text3),
                    filename,
                )
                return text3
            self._logger.debug("[KNOWLEDGE] pdfplumber empty for %r", filename)
        except Exception as exc:
            self._logger.warning(
                "[KNOWLEDGE] pdfplumber failed for %r: %s: %s",
                filename,
                type(exc).__name__,
                exc,
            )

        # --- All extractors exhausted (likely a scanned / image-only PDF) ---
        self._logger.warning(
            "[KNOWLEDGE] All PDF extractors returned empty for %r "
            "(scanned/image-only PDF?). Indexing with placeholder so blob is not orphaned.",
            filename,
        )
        print(
            f"[KNOWLEDGE] WARNING: no extractable text in PDF {filename!r}; "
            "indexing with '[No extractable text]' placeholder"
        )
        return "[No extractable text]"


__all__ = [
    "KnowledgeIngestionService",
    "KnowledgeSearchIndex",
]
